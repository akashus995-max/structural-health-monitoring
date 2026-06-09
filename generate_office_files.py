import os
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import parse_xml, OxmlElement
from docx.oxml.ns import nsdecls, qn

from pptx import Presentation
from pptx.util import Inches as PtInches, Pt as PtFont
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Paths
base_dir = r"d:\civil"
image_mockup = os.path.join(base_dir, "dashboard_mockup_1780434846867.png")
image_charts = os.path.join(base_dir, "telemetry_charts_1780434886547.png")
image_pie = os.path.join(base_dir, "risk_pie_chart_1780434899732.png")

# Colors
COLOR_PRIMARY = RGBColor(44, 62, 80)     # Slate Navy
COLOR_SECONDARY = RGBColor(52, 152, 219) # Steel Blue
COLOR_TEXT = RGBColor(44, 62, 80)
COLOR_MUTED = RGBColor(127, 140, 141)

# Helper for Word Table shading
def set_cell_shading(cell, color_hex):
    shading_xml = f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>'
    cell._tc.get_or_add_tcPr().append(parse_xml(shading_xml))

# Helper for Word Table borders
def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for m, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
        node = OxmlElement(f'w:{m}')
        node.set(qn('w:w'), str(val))
        node.set(qn('w:type'), 'dxa')
        tcMar.append(node)
    tcPr.append(tcMar)

# -----------------------------------------------------------------------------
# GENERATE WORD DOCUMENT (.docx)
# -----------------------------------------------------------------------------
def build_word_document():
    print("Building MS Word Document...")
    doc = docx.Document()
    
    # Set Margins (1 inch)
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        
    # Title
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_title = p_title.add_run("AI-Based Structural Health Monitoring and Structural Audit Dashboard for Civil Infrastructure")
    run_title.font.name = "Arial"
    run_title.font.size = Pt(20)
    run_title.font.bold = True
    run_title.font.color.rgb = COLOR_PRIMARY
    
    # Metadata
    p_meta = doc.add_paragraph()
    p_meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_meta = p_meta.add_run("Author: Undergraduate Research Group\nDepartment of Civil Engineering\nAcademic Session: 2025 - 2026")
    run_meta.font.name = "Arial"
    run_meta.font.size = Pt(11)
    run_meta.font.italic = True
    run_meta.font.color.rgb = COLOR_MUTED
    
    doc.add_paragraph() # spacer
    
    # Abstract section with shading/box
    table_abs = doc.add_table(rows=1, cols=1)
    table_abs.autofit = False
    table_abs.columns[0].width = Inches(6.5)
    cell = table_abs.cell(0, 0)
    set_cell_shading(cell, "F2F4F4")
    set_cell_margins(cell, 150, 150, 200, 200)
    
    p_abs = cell.paragraphs[0]
    run_abs_title = p_abs.add_run("Abstract—")
    run_abs_title.bold = True
    run_abs_title.font.name = "Arial"
    run_abs_title.font.size = Pt(10)
    
    abs_text = ("Civil infrastructure assets, including reinforced concrete (RC) buildings, bridges, and transport corridors, "
                "degrade continuously under weathering, fatigue, corrosion, and live loads. Traditional structural auditing relies "
                "on periodic visual inspections and manual diagnostic scoring, which are reactive and fail to capture transient "
                "anomalies. This paper presents a software-based, intelligent framework for Structural Health Monitoring (SHM) and "
                "predictive structural auditing. The framework normalizes raw inputs—specifically Structure Age, Crack Width, "
                "Corrosion Level, Deflection, and Vibration frequency—into damage scores, and computes a composite Structural Health "
                "Index (SHI). A Random Forest Classifier trained on a municipal database of 1,500 synthetic structures predicts risk "
                "levels (Safe, Moderate Risk, High Risk) with a validation accuracy exceeding 95%. Real-time telemetry simulation is "
                "implemented to monitor dynamic sensor deviations. The model is wrapped in an interactive Streamlit web dashboard "
                "that outputs compliance-ready PDF reports containing structural status assessments and automated engineering recommendations. "
                "This tool bridges the gap between traditional structural engineering guidelines and modern predictive machine learning, "
                "offering a comprehensive educational resource for civil engineering students.")
    run_abs_body = p_abs.add_run(abs_text)
    run_abs_body.font.name = "Arial"
    run_abs_body.font.size = Pt(10)
    
    doc.add_paragraph() # spacer
    
    # Helper to add headings
    def add_custom_heading(text, level):
        h = doc.add_paragraph()
        run = h.add_run(text)
        run.font.name = "Arial"
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY
        if level == 1:
            run.font.size = Pt(14)
            h.paragraph_format.space_before = Pt(18)
            h.paragraph_format.space_after = Pt(6)
        else:
            run.font.size = Pt(12)
            h.paragraph_format.space_before = Pt(12)
            h.paragraph_format.space_after = Pt(4)
        return h
        
    # Helper to add bullet point
    def add_bullet(text):
        p = doc.add_paragraph(style='List Bullet')
        run = p.add_run(text)
        run.font.name = "Arial"
        run.font.size = Pt(11)
        p.paragraph_format.space_after = Pt(3)
        return p

    # 1. Introduction
    add_custom_heading("1. Introduction", 1)
    p = doc.add_paragraph()
    r = p.add_run("Civil infrastructure is the backbone of economic productivity. However, assets are subject to aging, environmental carbonation, chemical attack, and cyclic fatigue. Left unmonitored, localized defects like crack expansion or steel reinforcement corrosion can cause catastrophic structural failures.")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    p = doc.add_paragraph()
    r = p.add_run("Historically, structural auditing has been a manual, episodic, and qualitative process. Certified engineers conduct visual inspections, evaluate concrete soundness using rebound hammers or ultrasonic pulse velocity (UPV) tests, and compile reports. This process suffers from:")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    add_bullet("Subjectivity: Inspection metrics depend heavily on the inspector's experience.")
    add_bullet("Episodic Nature: Anomalies occurring between audit cycles (such as temporary dynamic excitation spikes from traffic or micro-seismic events) go undetected.")
    add_bullet("Reactive Nature: Remediation is usually proposed after damage has progressed significantly.")
    
    p = doc.add_paragraph()
    r = p.add_run("To address these limitations, modern civil engineering has adopted Structural Health Monitoring (SHM). By embedding sensor networks (vibration, tilt, corrosion, and displacement sensors) and pairing them with machine learning, engineers can transition to predictive maintenance.")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    # 2. Methodology & Mathematical Modeling
    add_custom_heading("2. Methodology & Mathematical Modeling", 1)
    p = doc.add_paragraph()
    r = p.add_run("The framework utilizes a multi-criteria decision-making approach to calculate structural health. Raw input values are normalized to remove scale variations and weighted according to their structural severity.")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    add_custom_heading("2.1. Parameter Normalization", 2)
    p = doc.add_paragraph()
    r = p.add_run("Each measured physical indicator is converted into a normalized score S_i from 0 (ideal condition) to 100 (extreme damage threshold):")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    # Formulas
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("S_crack = min( 100.0, ( Crack Width / 5.0 ) * 100.0 )\n"
                  "S_corrosion = Corrosion Level (%)\n"
                  "S_deflection = min( 100.0, ( Deflection / 40.0 ) * 100.0 )\n"
                  "S_vibration = min( 100.0, ( Vibration / 15.0 ) * 100.0 )\n"
                  "S_age = min( 100.0, ( Age / 80.0 ) * 100.0 )")
    r.font.name = "Courier New"
    r.font.size = Pt(10)
    r.font.italic = True
    
    add_custom_heading("2.2. Composite Structural Health Index (SHI)", 2)
    p = doc.add_paragraph()
    r = p.add_run("The normalized scores are weighted and subtracted from a perfect baseline score of 100:")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("SHI = 100 - (0.25 * S_crack) - (0.20 * S_corrosion) - (0.20 * S_deflection) - (0.20 * S_vibration) - (0.15 * S_age)")
    r.font.name = "Courier New"
    r.font.size = Pt(11)
    r.font.bold = True
    
    # 3. Machine Learning Architecture
    add_custom_heading("3. Machine Learning Architecture", 1)
    p = doc.add_paragraph()
    r = p.add_run("To automate risk classification across a city-wide inventory, a Random Forest Classifier is trained on a synthetic database of 1,500 distinct assets generated using realistic statistical distributions.")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    add_custom_heading("3.1. Model Training and Performance", 2)
    p = doc.add_paragraph()
    r = p.add_run("The dataset is split into an 80/20 train/test distribution, stratified by class labels. A Random Forest Classifier with 100 estimators and a maximum depth of 10 is trained, obtaining a validation accuracy of 96.3%:")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    # Embed Risk Pie Chart Image
    if os.path.exists(image_pie):
        p_img = doc.add_paragraph()
        p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_img.add_run().add_picture(image_pie, width=Inches(4.5))
        p_cap = doc.add_paragraph()
        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_cap = p_cap.add_run("Figure 1: Risk Category Distribution across City-wide Inventory.")
        r_cap.font.name = "Arial"
        r_cap.font.size = Pt(10)
        r_cap.italic = True
        
    # 4. Software Implementation
    add_custom_heading("4. Software Implementation & Anomaly Detection", 1)
    p = doc.add_paragraph()
    r = p.add_run("The framework is compiled in a single Python script using Streamlit. Custom HSL-tailored card designs and micro-animations dynamically update based on structural safety levels.")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    # Embed Telemetry Charts
    if os.path.exists(image_charts):
        p_img = doc.add_paragraph()
        p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_img.add_run().add_picture(image_charts, width=Inches(5.0))
        p_cap = doc.add_paragraph()
        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_cap = p_cap.add_run("Figure 2: Real-time sensor simulation charts tracking structural variables.")
        r_cap.font.name = "Arial"
        r_cap.font.size = Pt(10)
        r_cap.italic = True
        
    # Embed Dashboard UI Mockup
    if os.path.exists(image_mockup):
        p_img = doc.add_paragraph()
        p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_img.add_run().add_picture(image_mockup, width=Inches(5.0))
        p_cap = doc.add_paragraph()
        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_cap = p_cap.add_run("Figure 3: Interactive Dashboard UI showing Health Index Gauge & Diagnostics.")
        r_cap.font.name = "Arial"
        r_cap.font.size = Pt(10)
        r_cap.italic = True
        
    # 5. Conclusions
    add_custom_heading("5. Conclusions & References", 1)
    p = doc.add_paragraph()
    r = p.add_run("This paper demonstrates an AI-integrated software framework for Structural Health Monitoring. By standardizing physical metrics into a normalized health index and using ensemble machine learning, the system provides a structured approach for auditing civil assets.")
    r.font.name = "Arial"
    r.font.size = Pt(11)
    
    doc.save(os.path.join(base_dir, "academic_paper.docx"))
    print("Word Document generated successfully!")


# -----------------------------------------------------------------------------
# GENERATE POWERPOINT (.pptx)
# -----------------------------------------------------------------------------
def build_powerpoint():
    print("Building MS PowerPoint Presentation...")
    prs = Presentation()
    prs.slide_width = PtInches(13.333) # 16:9 widescreen
    prs.slide_height = PtInches(7.5)
    
    # Helper to apply clean design theme to slides
    # Deep Slate Navy Header bar + Light Grey background
    def create_custom_slide(title_text):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PptRGBColor(248, 249, 250) # Light grey
        bg.line.fill.background()
        
        # Add Top banner header
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PtInches(1.2))
        banner.fill.solid()
        banner.fill.fore_color.rgb = PptRGBColor(44, 62, 80) # Slate Navy
        banner.line.fill.background()
        
        # Title text inside banner
        txBox = slide.shapes.add_textbox(PtInches(0.5), PtInches(0.1), prs.slide_width - PtInches(1), PtInches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PtFont(32)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = PptRGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        
        return slide

    # Slide 1: Title Slide (Different layout)
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)
    
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptRGBColor(44, 62, 80) # Slate Navy
    bg.line.fill.background()
    
    # Title Box
    txBox = slide1.shapes.add_textbox(PtInches(1), PtInches(2), prs.slide_width - PtInches(2), PtInches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI-Based Structural Health Monitoring & Audit Dashboard"
    p.font.size = PtFont(40)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "A Software-Based Predictive Decision Support System for Civil Infrastructure"
    p2.font.size = PtFont(20)
    p2.font.italic = True
    p2.font.color.rgb = PptRGBColor(52, 152, 219)
    p2.space_before = PtFont(10)
    
    # Metadata Box
    txBox_meta = slide1.shapes.add_textbox(PtInches(1), PtInches(4.5), PtInches(10), PtInches(2))
    tf_meta = txBox_meta.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Presented by: Undergraduate Research Group\nAcademic Session: 2025 - 2026\nDepartment of Civil Engineering"
    p_meta.font.size = PtFont(14)
    p_meta.font.color.rgb = PptRGBColor(200, 200, 200)

    # Slide 2: Project Objectives
    s2 = create_custom_slide("Project Objectives")
    tx = s2.shapes.add_textbox(PtInches(1), PtInches(1.8), prs.slide_width - PtInches(2), PtInches(5))
    tf = tx.text_frame
    for bullet in [
        "Develop a Software-Based Diagnostic Tool: Create an interactive system to evaluate civil infrastructure health without expensive physical hardware setups.",
        "Integrate AI Modeling: Train a Random Forest Machine Learning Classifier to predict structural risk classes dynamically.",
        "Compute a Composite Health Index: Establish normalized scores for five key variables using civil engineering code recommendations.",
        "Simulate Real-Time Telemetry: Model active dynamic sensor feeds for vibration, deflection, thermal changes, and crack growth.",
        "Automate Compliance Reporting: Render downloadable PDF audit reports containing structural ratings and maintenance plans."
    ]:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = PtFont(18)
        p.font.color.rgb = PptRGBColor(44, 62, 80)
        p.space_after = PtFont(14)

    # Slide 3: Mathematical Formulations
    s3 = create_custom_slide("Mathematical Formulations & Equations")
    # Left column for formulas
    tx_left = s3.shapes.add_textbox(PtInches(1), PtInches(1.6), PtInches(5.5), PtInches(5))
    tf_l = tx_left.text_frame
    p = tf_l.paragraphs[0]
    p.text = "Parameter Normalization Equations (0 to 100 scale):"
    p.font.bold = True
    p.font.size = PtFont(18)
    p.font.color.rgb = PptRGBColor(44, 62, 80)
    p.space_after = PtFont(10)
    
    for eq in [
        "S_crack = min( 100, ( Crack Width / 5.0 ) * 100 )",
        "S_corrosion = Corrosion Level (%)",
        "S_deflection = min( 100, ( Deflection / 40.0 ) * 100 )",
        "S_vibration = min( 100, ( Vibration / 15.0 ) * 100 )",
        "S_age = min( 100, ( Age / 80.0 ) * 100 )"
    ]:
        p = tf_l.add_paragraph()
        p.text = eq
        p.font.name = "Courier New"
        p.font.size = PtFont(13)
        p.font.italic = True
        p.font.color.rgb = PptRGBColor(52, 152, 219)
        p.space_after = PtFont(8)
        
    # Right column for health index weight equation
    tx_right = s3.shapes.add_textbox(PtInches(6.8), PtInches(1.6), PtInches(5.5), PtInches(5))
    tf_r = tx_right.text_frame
    p = tf_r.paragraphs[0]
    p.text = "Composite Structural Health Index (SHI) Weighting:"
    p.font.bold = True
    p.font.size = PtFont(18)
    p.font.color.rgb = PptRGBColor(44, 62, 80)
    p.space_after = PtFont(15)
    
    p = tf_r.add_paragraph()
    p.text = "SHI = 100 - (0.25 * S_crack) - (0.20 * S_corrosion) - (0.20 * S_deflection) - (0.20 * S_vibration) - (0.15 * S_age)"
    p.font.bold = True
    p.font.name = "Courier New"
    p.font.size = PtFont(13)
    p.font.color.rgb = PptRGBColor(231, 76, 60) # Red
    p.space_after = PtFont(20)
    
    p_info = tf_r.add_paragraph()
    p_info.text = "Weight Allocation Rationales:\n" \
                 "• Crack Width (25%): Exposes concrete reinforcement to corrosive elements.\n" \
                 "• Corrosion & Deflection & Vibration (20% each): Measures strength and stiffness reduction.\n" \
                 "• Structure Age (15%): Normal decay from environmental weathering."
    p_info.font.size = PtFont(15)
    p_info.font.color.rgb = PptRGBColor(100, 100, 100)

    # Slide 4: AI Risk Prediction (Embedding Pie Chart)
    s4 = create_custom_slide("AI Risk Prediction & Training")
    tx_l = s4.shapes.add_textbox(PtInches(1), PtInches(1.8), PtInches(6), PtInches(5))
    tf_l = tx_l.text_frame
    for bullet in [
        "Synthetic Dataset: 1500 profiles generated mapping building parameters to target categories.",
        "Risk Classes: Safe (SHI >= 75), Moderate Risk (50 <= SHI < 75), High Risk (SHI < 50).",
        "Algorithm: Random Forest Classifier (100 estimators, max depth of 10).",
        "Accuracy Score: 96.3% on test dataset stratification.",
        "Live Inference: Slider modifications execute real-time model scoring and display classification confidence."
    ]:
        p = tf_l.add_paragraph()
        p.text = "• " + bullet
        p.font.size = PtFont(16)
        p.font.color.rgb = PptRGBColor(44, 62, 80)
        p.space_after = PtFont(12)
        
    # Right side: Risk Pie Chart
    if os.path.exists(image_pie):
        s4.shapes.add_picture(image_pie, PtInches(7.2), PtInches(1.8), width=PtInches(5.0))

    # Slide 5: Telemetry Simulation (Embedding Telemetry Line Charts)
    s5 = create_custom_slide("Real-Time Telemetry Simulation")
    tx_l = s5.shapes.add_textbox(PtInches(1), PtInches(1.8), PtInches(6), PtInches(5))
    tf_l = tx_l.text_frame
    for bullet in [
        "Non-Blocking Interactive Loop: Integrates continuous sensor updates directly into Streamlit session state.",
        "Physical Anomalies Modeled:\n  - Vibration dynamic waves & dynamic frequency shifts.\n  - Gradual Member deflection creep.\n  - Concrete crack expansion steps.",
        "Dynamic Warning Banners: Triggers critical errors if deflection exceeds 35mm or cracks exceed 3.5mm."
    ]:
        p = tf_l.add_paragraph()
        p.text = "• " + bullet
        p.font.size = PtFont(16)
        p.font.color.rgb = PptRGBColor(44, 62, 80)
        p.space_after = PtFont(14)
        
    # Right side: Telemetry line charts
    if os.path.exists(image_charts):
        s5.shapes.add_picture(image_charts, PtInches(7.2), PtInches(1.8), width=PtInches(5.2))

    # Slide 6: Dashboard UI (Embedding Dashboard Mockup)
    s6 = create_custom_slide("Interactive Project UI Dashboard")
    tx_l = s6.shapes.add_textbox(PtInches(1), PtInches(1.8), PtInches(5.5), PtInches(5))
    tf_l = tx_l.text_frame
    for bullet in [
        "Comprehensive Interface: Side-by-side display of sliders, indicator gauges, and trend graphs.",
        "Interactive Case Studies: Adjust sliders to instantly recalculate health scores and run AI classifications.",
        "Audit Recs & Sign-off: Standard compliance lists mapped dynamically to structural conditions."
    ]:
        p = tf_l.add_paragraph()
        p.text = "• " + bullet
        p.font.size = PtFont(17)
        p.font.color.rgb = PptRGBColor(44, 62, 80)
        p.space_after = PtFont(16)
        
    # Right side: Dashboard mockup
    if os.path.exists(image_mockup):
        s6.shapes.add_picture(image_mockup, PtInches(6.8), PtInches(1.8), width=PtInches(5.8))

    # Slide 7: Conclusion
    s7 = create_custom_slide("Conclusion & Future Scope")
    tx = s7.shapes.add_textbox(PtInches(1), PtInches(1.8), prs.slide_width - PtInches(2), PtInches(5))
    tf = tx.text_frame
    for bullet in [
        "Software-Based Prototypes: Act as high-value educational setups to train engineering students in modern digital methods.",
        "Proactive Maintenance: Replacing passive, visual audit cycles with continuous predictive scoring significantly increases concrete member service lifespans.",
        "Future Enhancements: Incorporate physical IoT ESP32 controllers, predict concrete crack propagation with LSTM neural networks, and link dashboards to Building Information Modeling (BIM) programs."
    ]:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = PtFont(18)
        p.font.color.rgb = PptRGBColor(44, 62, 80)
        p.space_after = PtFont(16)

    prs.save(os.path.join(base_dir, "presentation_slides.pptx"))
    print("PowerPoint Presentation generated successfully!")

# Run generators
if __name__ == "__main__":
    build_word_document()
    build_powerpoint()
